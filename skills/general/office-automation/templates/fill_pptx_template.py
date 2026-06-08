"""
Template: Fill a PowerPoint presentation from research content.

WORKFLOW:
1. Replace the source paths with actual file paths
2. Update the content strings to match your document
3. Run: python fill_pptx_template.py

PITFALL: The dst must be in a DIFFERENT directory than src.
If the source .pptx is open in PowerPoint, writing to the same
directory will fail with PermissionError.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil

# Paths — dst MUST be different directory than src
src = r'C:\path\to\template.pptx'
dst = r'C:\path\to\output\filled.pptx'

shutil.copy2(src, dst)
prs = Presentation(dst)

black = RGBColor(0x33, 0x33, 0x33)

def add_tb(slide, left, top, width, height, text, font_size=11, bold=False, color=None, align=None):
    """Add a textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size: run.font.size = Pt(font_size)
    if bold: run.font.bold = True
    if color: run.font.color.rgb = color
    if align is not None: p.alignment = align
    return txBox

def clear_and_set(shape, old_sub, new_text, fs=None, bld=None):
    """Find old_sub in shape text and replace with new_text."""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        if old_sub in para.text.strip():
            for run in para.runs:
                run.text = ''
            if para.runs:
                para.runs[0].text = new_text
                if fs: para.runs[0].font.size = Pt(fs)
                if bld is not None: para.runs[0].font.bold = bld
                return True
            else:
                r = para.add_run()
                r.text = new_text
                if fs: r.font.size = Pt(fs)
                if bld is not None: r.font.bold = bld
                return True
    return False

# ===== INSPECT TEMPLATE FIRST =====
# Uncomment to see what shapes exist:
# for i, slide in enumerate(prs.slides):
#     print(f'Slide {i+1}:')
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for para in shape.text_frame.paragraphs:
#                 if para.text.strip():
#                     print(f'  {shape.name}: "{para.text[:80]}"')

# ===== FILL SLIDES =====
# Example for slide 1:
s1 = prs.slides[0]
for sh in s1.shapes:
    clear_and_set(sh, 'Placeholder Title', 'Your Title Here', fs=16, bld=True)
    clear_and_set(sh, 'Placeholder Date', 'May 2026')

# Add content textbox
add_tb(s1, Inches(0.8), Inches(1.8), Inches(8.2), Inches(3.5),
    'Your content here. Max 6 bullets per slide.',
    font_size=12, color=black)

prs.save(dst)
print(f'LISTO: {dst}')
